import json

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import openpyxl
from openpyxl import Workbook
import datetime

CN_PPTFILE='D:\\cn_v2.pptx'
EN_PPTFILE='D:\\en_v2.pptx'
EXCELFILE="translate_word.xlsx"
TYPE_TEXT='text'
TYPE_TABLE='tablet'
SHEET_NAME_TEXT='text_cn_en'
SHEET_NAME_TABLE='table_cn_en'

def judge_pure_english(keyword):
    return all(ord(c) < 128 for c in keyword)

def appendToList(raw_text, textList):
    print(raw_text)
    if judge_pure_english(raw_text):
        print('pure english,ignore,continue... ...')
        return
    textList.append(raw_text)

def appendTableToList(raw_text, textList):
    textList.append(raw_text)

def deleteDupElemFromList(textList):
    return list(set(textList))
    #return sorted(set(textList), key = textList.index)

#get all text from ppt file
def ppt_catch_format_text(filename):
    """
    抓取PPT的内容，按段落返回
    其中 filename 是PPT文件的路径
    """
    prs = Presentation(filename)
    txtList = []
    for x in range(len(prs.slides)):
        # ---Only on text-boxes outside group elements---
        for shape in prs.slides[x].shapes:
            print(shape)
            if hasattr(shape, "text"):
                print('hasattr text in shape%s'%(shape))
                raw_text = shape.text.encode('utf-8').strip().decode()
                appendToList(raw_text,txtList)
                #for paragraph in shape.text_frame.paragraphs:
                #    print('----from paragraph-----')
                #    appendToList(paragraph.text, txtList)
        # ---Only operate on group shapes---
        group_shapes = [shp for shp in prs.slides[x].shapes 
                        if shp.shape_type ==MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            print(group_shape)
            for shape in group_shape.shapes:
                print(shape)
                if shape.has_text_frame:
                    print('has_text_frame in shape%s'%(shape))
                    #row_text=shape.text_framew.text
                    raw_text = shape.text.encode('utf-8').strip().decode()
                    appendToList(raw_text, txtList)
                    
                    #for paragraph in shape.text_frame.paragraphs:
                    #    print('----from paragraph-----')
                    #    appendToList(paragraph.text, txtList)
    
    prs.save(filename)
    return txtList

def tranvel_table(table, txtList):
    rows=len(table.rows)
    columns=len(table.columns)
    for i in range(rows):
        for j in range(columns):
            raw_text=table.cell(i,j).text
            print('tranvel_table text:%s'%(raw_text))
            appendToList(raw_text,txtList)
            #table.cell(i,j).text='A'+text

def tranvel_table_replace(table, text_list_cn, text_list_en):
    rows=len(table.rows)
    columns=len(table.columns)
    for i in range(rows):
        for j in range(columns):
            raw_text=table.cell(i,j).text
            print('tranvel_table text:%s'%(raw_text))
            raw_text = shape.text.encode('utf-8').strip().decode()
            targetText=getTranslate(raw_text,text_list_cn,text_list_en)
            print('targetText:%s'%targetText)
            if targetText != None:
                table.cell(i,j).text=targetText
                print('replace ok')

def ppt_catch_format_text_from_table(filename):
    txtList = []
    prs = Presentation(filename)

    for x in range(len(prs.slides)):
        # ---Only on table outside group elements---
        for shape in prs.slides[x].shapes:
            print(shape)
            if shape.has_table:
                print('1has_table in shape%s'%(shape))
                table = shape.table
                tranvel_table(table, txtList)
                #appendTableToList(row_text,txtList)
                #for paragraph in shape.text_frame.paragraphs:
                #    print('----from paragraph-----')
                #    appendToList(paragraph.text, txtList)
        # ---Only operate on group shapes---
        group_shapes = [shp for shp in prs.slides[x].shapes 
                        if shp.shape_type ==MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            print(group_shape)
            for shape in group_shape.shapes:
                print(shape)
                if shape.has_table:
                    print('g has_table in shape%s'%(shape))
                    #row_text=shape.text_framew.text
                    row_text = shape.table
                    appendTableToList(row_text, txtList)
                    
                    #for paragraph in shape.text_frame.paragraphs:
                    #    print('----from paragraph-----')
                    #    appendToList(paragraph.text, txtList)

    prs.save(filename)
    return txtList


def ppt_catch_format_text_from_datasource(filename):
   txtList = []
   return txtList

#write word list to excel file
def write_excel(filename, txt_list_cn, type):
    text_list = deleteDupElemFromList(txt_list_cn)

    if type==TYPE_TEXT:
        sheetname=SHEET_NAME_TEXT
    elif type==TYPE_TABLE:
        sheetname=SHEET_NAME_TABLE
    else:
        return
    #wb = Workbook()
    wb = openpyxl.load_workbook(filename)
    word_sheet = wb.create_sheet(sheetname)
    #rowindex colindex startswith 1 in table
    rowindex=1
    for wordtxt in text_list:
        print('s%se'%(wordtxt))
        #wordtxt.replace('\\','\\\\')
        if judge_pure_english(wordtxt):
            print('pure english,ignore,continue... ...')
            continue
        word_sheet.cell(rowindex,1).value=wordtxt
        rowindex=rowindex+1
        #wordfile.write(wordtxt)
        #wordfile.write(",")
    wb.save(filename)

def read_excel(filename, type):
    text_list_cn = []
    text_list_en = []
    wb = openpyxl.load_workbook(filename)
    #ws = wb.active
    
    if type==TYPE_TEXT:
        sheetname=SHEET_NAME_TEXT
    elif type==TYPE_TABLE:
        sheetname=SHEET_NAME_TABLE
    else:
        return
    word_sheet = wb[sheetname]
    #print(list(word_sheet.columns)[0])
    colA=word_sheet['A']

    #print(colA)
    #print(type(colA))
    #print(colA.__class__)
    #print(dir(colA))
    colB=word_sheet['B']
    for cell in colA:
        text_list_cn.append(cell.value)
    for cell in colB:
        text_list_en.append(cell.value)
    print(text_list_cn)
    print(text_list_en)        

    return text_list_cn,text_list_en


#获得cn对应的翻译en
def getTranslate(text,text_list_cn,text_list_en):
    for index in range(len(text_list_cn)):
        #print('search:%s index:%s value:%s'%(text,index,text_list_cn[index]))
        if text_list_cn[index].strip() == text.strip():
            return text_list_en[index]

def replace_ppt_text(filename, text_list_cn, text_list_en, targetfilename):
    prs = Presentation(filename)
    for x in range(len(prs.slides)):
        #txt_oa[x] = []
        # ---Only on text-boxes outside group elements---
        for shape in prs.slides[x].shapes:
            print(shape)
            if hasattr(shape, "text"):
                print(shape.text)
                raw_text = shape.text.encode('utf-8').strip().decode()
                print('Sshape.textInPpt:%sE'%(shape.text))
                print('Srow_textInPpt:%sE'%(raw_text))
                targetText=getTranslate(raw_text,text_list_cn,text_list_en)
                print('targetText:%s'%targetText)
                if targetText != None:
                    shape.text=targetText
                    print('replace ok')

        # ---Only operate on group shapes---
        group_shapes = [shp for shp in prs.slides[x].shapes 
                        if shp.shape_type ==MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    raw_text = shape.text.encode('utf-8').strip().decode()
                    print('Sshape.textInPpt:%sE'%(shape.text))
                    print('Srow_textInPpt:%sE'%(raw_text))
                    targetText=getTranslate(raw_text,text_list_cn,text_list_en)
                    print('targetText:%s'%targetText)
                    if targetText != None:
                        shape.text=targetText
                        print('replace ok')
    
    prs.save(targetfilename)

def replace_ppt_table(filename, text_list_cn, text_list_en, targetfilename):
    prs = Presentation(filename)
    for x in range(len(prs.slides)):
        #txt_oa[x] = []
        # ---Only on text-boxes outside group elements---
        for shape in prs.slides[x].shapes:
            if shape.has_table:
                print('1has_table in shape%s'%(shape))
                table = shape.table
                tranvel_table_replace(table, text_list_cn, text_list_en)

    prs.save(targetfilename)


    
def python_help():
    print('---------------get cn word from pptx-----------------------')
    text_list = ppt_catch_format_text('E:\D\sec.pptx')
    print('---------------write word to xlsx-----------------------')
    write_excel("word.xlsx",text_list)
    #text_list = json.dumps(text_list, ensure_ascii=False, indent=4).replace("\\n","")
    print('--------------------------------------')
    print(text_list)
    print(text_list.__class__)
    #print(text_list.__base__)
    print(dir(text_list))

    print(type(text_list))
    print(type(type(text_list)))

    print('--------------------------------------')
    #print(text_list.object)
    #print(text_list.__str__)
    #wordfile.write("\n")

CMD1='READ text'
CMD2='READ table'
CMD3='WRITE text'
CMD4='WRITE table'
def main():
    print('---------------enter main()-----------------------')
    cmd=1
    print('---------------enter CMD:%d-----------------------'%cmd)
    if cmd==1:
        text_list=ppt_catch_format_text(CN_PPTFILE)
        print(text_list)
        write_excel(EXCELFILE,text_list,TYPE_TEXT)
    elif cmd==2:
        text_list=ppt_catch_format_text_from_table(CN_PPTFILE)
        print(text_list)
        write_excel(EXCELFILE,text_list,TYPE_TABLE)
    elif cmd==3:
        #maybe excel format changed, cannot compare
        text_list_cn,text_list_en=read_excel(EXCELFILE,TYPE_TEXT)
        print('---------------after read_excel-----------------------')
        print(text_list_cn)
        print(text_list_en)
        print('---------------enter replace_ppt-----------------------')
        replace_ppt_text(CN_PPTFILE,text_list_cn, text_list_en, EN_PPTFILE)
    elif cmd==4:
        text_list_cn,text_list_en=read_excel(EXCELFILE,TYPE_TABLE)
        print('---------------after read_excel-----------------------')
        print(text_list_cn)
        print(text_list_en)
        print('---------------enter replace_ppt-----------------------')
        replace_ppt_table(CN_PPTFILE,text_list_cn, text_list_en, EN_PPTFILE)

    
main()
#AF_INET域 AF_UNIX